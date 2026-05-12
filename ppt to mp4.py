from pptx import Presentation
from pptx.util import Inches
import os
import subprocess

# =========================
# FILES
# =========================
ffmpeg_path = r"C:\ffmpeg\bin\ffmpeg.exe"
def add_footer(ppt_input,ppt_output,bottom_image):
    # ppt_input = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier.pptx"
    # ppt_output = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier_out.pptx"
    #
    # bottom_image = r"C:\Users\Shubham Aggarwal\Downloads\Banner.jpg"

    # =========================
    # LOAD PPT
    # =========================

    prs = Presentation(ppt_input)

    # =========================
    # ADD IMAGE TO ALL SLIDES
    # =========================

    for slide in prs.slides:

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Bottom image dimensions
        img_height = Inches(0.4)

        slide.shapes.add_picture(
            bottom_image,
            0,
            slide_height - img_height,
            width=slide_width,
            height=img_height
        )

    # =========================
    # SAVE PPT
    # =========================

    prs.save(ppt_output)

    print("Modified PPT saved.")

def export_as_video(ppt_path,video_path):
    import win32com.client
    import time

    # ppt_path = r"C:\path\output_modified.pptx"
    # video_path = r"C:\path\final_video.mp4"

    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(ppt_path)

    # Export as MP4
    presentation.CreateVideo(
        video_path,
        True,  # Use Timings
        10,  # Seconds per slide
        720,  # Vertical resolution
        30,  # FPS
        85  # Quality
    )

    # Wait until video creation finishes
    while presentation.CreateVideoStatus == 1:
        print("Creating video...")
        time.sleep(5)

    presentation.Close()
    powerpoint.Quit()

    print("Video exported.")

def merge_music_with_video(video,music,output):
    import os
    print(video)

    # video = "final_video.mp4"
    # music = "background.mp3"
    # output = "final_with_music.mp4"

    # cmd = f'ffmpeg -i "{video}" -i "{music}" -c:v copy -c:a aac -shortest "{output}"'
    subprocess.run([
        ffmpeg_path,
        "-i", video,
        "-i", music,
        "-c:v", "copy",
        "-c:a", "aac",
        "-shortest",
        output
    ])

    # os.system(cmd)

ppt_input = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier.pptx"
ppt_output = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier_out.pptx"
bottom_image = r"C:\Users\Shubham Aggarwal\Downloads\Banner.jpg"
video_path = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier_out.mp4"
video_name = "AIIMS GOVT Awantipora_Dossier"
music = r"C:\Users\Shubham Aggarwal\Downloads\WhatsApp Audio 2026-05-12 at 4.13.05 PM.mpeg"
video_output = r"C:\Users\Shubham Aggarwal\Downloads\AIIMS GOVT Awantipora_Dossier_out_video.mp4"

add_footer(ppt_input,ppt_output,bottom_image)
export_as_video(ppt_output,video_path)
merge_music_with_video(video_path,music,video_output)